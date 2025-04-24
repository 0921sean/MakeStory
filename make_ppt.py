import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

# 원하는 주차 리스트
target_weeks = [4, 5, 6]

# 이미지 폴더 경로
base_folder = "ppt_pictures"

# 저장할 최종 경로
output_root = "고문상_강의안"
os.makedirs(output_root, exist_ok=True)

# 슬라이드 크기: 16:9
slide_width = Inches(13.33)
slide_height = Inches(7.5)

for week in target_weeks:
    folder_name = f"{week}주차"
    week_path = os.path.join(base_folder, folder_name)

    if not os.path.exists(week_path):
        print(f"[경고] {folder_name} 폴더가 존재하지 않음")
        continue

    image_files = sorted([
        f for f in os.listdir(week_path)
        if f.lower().endswith(('.png', '.jpg', '.jpeg'))
    ])

    if not image_files:
        print(f"[알림] {folder_name} 폴더에 이미지가 없음")
        continue

    prs = Presentation()
    prs.slide_width = slide_width
    prs.slide_height = slide_height

    for img_file in image_files:
        img_path = os.path.join(week_path, img_file)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(img_path, 0, 0, width=slide_width, height=slide_height)

    # 파일명
    pptx_filename = f"고문상_{week}주차.pptx"
    pdf_filename = f"고문상_{week}주차.pdf"

    # 저장 경로
    pptx_path = os.path.join(output_root, pptx_filename)
    pdf_path = os.path.join(output_root, pdf_filename)

    # PPT 저장
    prs.save(pptx_path)
    print(f"[PPT 저장 완료] {pptx_path}")

    # PDF 저장 (Mac 기준: Preview나 Keynote로 가능하지만, 코드로는 다음처럼 PIL로 간단 처리)
    images = [Image.open(os.path.join(week_path, f)).convert("RGB") for f in image_files]
    if images:
        images[0].save(pdf_path, save_all=True, append_images=images[1:])
        print(f"[PDF 저장 완료] {pdf_path}")